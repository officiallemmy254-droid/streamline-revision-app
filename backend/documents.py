"""
Document router — PDF upload, listing, and deletion.

PDFs are processed into text chunks and stored in SQLite.
The raw PDF binary is discarded after extraction.
"""

from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, status, Form
from sqlalchemy.orm import Session
from pydantic import BaseModel
from datetime import datetime

import PyPDF2
import docx
import pptx
import io

from . import database, auth
from .retriever import chunk_text

router = APIRouter(prefix="/documents", tags=["Documents"])


# ── Schemas ─────────────────────────────────────────
class DocumentOut(BaseModel):
    id: int
    filename: str
    chunk_count: int
    category: str
    uploaded_at: str

    class Config:
        from_attributes = True


class DocumentListResponse(BaseModel):
    documents: list[DocumentOut]


class MessageResponse(BaseModel):
    message: str


# ── Dependency ──────────────────────────────────────
def get_db():
    db = database.SessionLocal()
    try:
        yield db
    finally:
        db.close()


# ── Helpers ─────────────────────────────────────────
def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract all text from a PDF file."""
    pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    text = ""
    for page in pdf_reader.pages:
        extracted = page.extract_text()
        if extracted:
            text += extracted + "\n"
    return text

def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract all text from a Word document."""
    doc = docx.Document(io.BytesIO(file_bytes))
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract all text from a PowerPoint presentation."""
    prs = pptx.Presentation(io.BytesIO(file_bytes))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


# ── Endpoints ───────────────────────────────────────
# Accepted document categories
VALID_CATEGORIES = {
    "TEACHING_DIPLOMA",   # legacy — Teaching Diploma materials
    "LEARNING_OUTCOMES",  # module/unit learning outcomes
    "NOTES",              # general handwritten or typed notes
    "LECTURE_NOTES",      # lecture slides and detailed notes
    "TEXTBOOK_CHAPTER",   # textbook excerpts or chapters
    "PAST_EXAM",          # past examination papers
}


@router.post("/upload", response_model=DocumentOut, status_code=status.HTTP_201_CREATED)
async def upload_document(
    file: UploadFile = File(...),
    category: str = Form("LECTURE_NOTES"),
    current_user: str = Depends(auth.get_current_user),
    db: Session = Depends(get_db),
):
    """Upload a PDF, extract text, chunk it, and store in the database."""
    # Validate category
    if category not in VALID_CATEGORIES:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid category '{category}'. Must be one of: {', '.join(sorted(VALID_CATEGORIES))}",
        )

    # Validate file type
    allowed_extensions = (".pdf", ".docx", ".pptx")
    if not file.filename or not file.filename.lower().endswith(allowed_extensions):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Only {', '.join(allowed_extensions)} files are accepted.",
        )

    # Read file
    file_bytes = await file.read()
    if len(file_bytes) > 10 * 1024 * 1024:  # 10MB limit
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="File size exceeds 10MB limit.",
        )

    # Extract text
    filename_lower = file.filename.lower()
    raw_text = ""
    try:
        if filename_lower.endswith(".pdf"):
            raw_text = extract_text_from_pdf(file_bytes)
        elif filename_lower.endswith(".docx"):
            raw_text = extract_text_from_docx(file_bytes)
        elif filename_lower.endswith(".pptx"):
            raw_text = extract_text_from_pptx(file_bytes)
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Error parsing file: {str(e)}"
        )

    if not raw_text.strip():
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Could not extract text from this file. It may be empty, scanned, or protected.",
        )

    # Get user from DB
    user = db.query(database.User).filter(database.User.username == current_user).first()
    if not user:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="User not found.")

    # Chunk text
    chunks = chunk_text(raw_text)

    # Create document record
    doc = database.Document(
        user_id=user.id,
        filename=file.filename,
        chunk_count=len(chunks),
        category=category,
    )
    db.add(doc)
    db.flush()  # Get the doc.id before adding chunks

    # Store chunks
    for idx, chunk_content in enumerate(chunks):
        chunk = database.DocumentChunk(
            document_id=doc.id,
            chunk_index=idx,
            content=chunk_content,
        )
        db.add(chunk)

    db.commit()
    db.refresh(doc)

    return DocumentOut(
        id=doc.id,
        filename=doc.filename,
        chunk_count=doc.chunk_count,
        category=doc.category,
        uploaded_at=doc.uploaded_at.isoformat(),
    )


@router.get("/", response_model=DocumentListResponse)
def list_documents(
    current_user: str = Depends(auth.get_current_user),
    db: Session = Depends(get_db),
):
    """List all documents for the authenticated user."""
    user = db.query(database.User).filter(database.User.username == current_user).first()
    if not user:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="User not found.")

    docs = db.query(database.Document).filter(database.Document.user_id == user.id).all()
    return DocumentListResponse(
        documents=[
            DocumentOut(
                id=d.id,
                filename=d.filename,
                chunk_count=d.chunk_count,
                category=d.category,
                uploaded_at=d.uploaded_at.isoformat(),
            )
            for d in docs
        ]
    )


@router.delete("/{document_id}", response_model=MessageResponse)
def delete_document(
    document_id: int,
    current_user: str = Depends(auth.get_current_user),
    db: Session = Depends(get_db),
):
    """Delete a document and all its chunks."""
    user = db.query(database.User).filter(database.User.username == current_user).first()
    if not user:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="User not found.")

    doc = db.query(database.Document).filter(
        database.Document.id == document_id,
        database.Document.user_id == user.id,
    ).first()

    if not doc:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Document not found.")

    db.delete(doc)  # Cascade deletes chunks
    db.commit()
    return {"message": f"Document '{doc.filename}' deleted."}
